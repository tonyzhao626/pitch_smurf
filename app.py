from flask import Flask, render_template, redirect, url_for
from flask_wtf import FlaskForm
from wtforms import StringField
import requests
from bs4 import BeautifulSoup
from tickerscrape import ticker
from tickerscrape import overview
from pptx import Presentation



app = Flask(__name__)
app.config['SECRET_KEY'] = 'SecretHere'


class AnotherForm(FlaskForm):
    company = StringField('Company Name : ')


#Works!
@app.route('/')
def default():
    return redirect(url_for('form3'))
    return "test"


#Works!
@app.route('/helloworld')
def hello_world():
   return render_template("index.html")



@app.route('/form', methods=['GET', 'POST'])
def form3():
    form3 = AnotherForm()
    if form3.validate_on_submit():
        if form3.company.data == "":
            return redirect(url_for('f404'))
        return redirect(url_for('fcomp',companyname = form3.company.data))
    return render_template('form3.html', form=form3)



@app.route('/company/<path:companyname>')
def fcomp(companyname):


    tickerurl = 'https://www.marketwatch.com/tools/quotes/lookup.asp?siteID=mktw&Lookup=' + companyname + '&Country=us&Type=All'
    content = requests.get(tickerurl)
    soup = BeautifulSoup(content.text, 'html.parser')


    name = companyname

    t = ticker(name)

    listofdata = overview(ticker(name)).split("---")

    root = Presentation()
    first_slide_layout = root.slide_layouts[0]
    slide = root.slides.add_slide(first_slide_layout)
    slide.shapes.title.text = name
    slide.placeholders[1].text = t + " -- Pitch Smurf"

    second_slide_layout = root.slide_layouts[1]
    slide2 = root.slides.add_slide(second_slide_layout)
    slide2.shapes.title.text = listofdata[0]
    slide2.placeholders[1].text = listofdata[1]

    third_slide_layout = root.slide_layouts[1]
    slide3 = root.slides.add_slide(third_slide_layout)
    slide3.shapes.title.text = "Key Shareholders"
    slide3.placeholders[1].text = listofdata[2] + " & " + listofdata[3]

    fourth_slide_layout = root.slide_layouts[1]
    slide4 = root.slides.add_slide(fourth_slide_layout)
    slide4.shapes.title.text = "Valuation"
    slide4.placeholders[1].text = "Market Cap: " + listofdata[4] + "P/E Ratio: " + listofdata[5] + "EPS: " + listofdata[6]

    fifth_slide_layout = root.slide_layouts[1]
    slide4 = root.slides.add_slide(fifth_slide_layout)
    slide4.shapes.title.text = "Predictions"
    slide4.placeholders[1].text = "Rating: " + listofdata[7] + "Expected Return on Equity in 6 months based on TTM: " + listofdata[8]

    #used ticker, sector, description, invest1, invest2

    root.save("Output.pptx")


    return overview(ticker(name))


@app.route('/404')
def f404():
    return render_template('404page.html')




if __name__ == '__main__':
    app.run(debug=True)
